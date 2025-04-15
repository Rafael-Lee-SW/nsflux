from fastapi import FastAPI, Request, UploadFile, File, Form, HTTPException
from fastapi.responses import JSONResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from typing import List, Optional
import os
import json
import datetime
import torch
import numpy as np
from pptx import Presentation
import PyPDF2
from pathlib import Path

# 임베딩 모델 로드
from transformers import AutoModel, AutoTokenizer
from utils import vectorize_content, normalize_text_vis

app = FastAPI()

# 정적 파일 및 템플릿 설정
BASE_DIR = Path(__file__).resolve().parent
app.mount("/static", StaticFiles(directory=str(BASE_DIR / "static")), name="static")
templates = Jinja2Templates(directory=str(BASE_DIR / "templates"))

# 데이터 경로 설정
DATA_PATH = BASE_DIR / "data" / "data.json"
DATA_PATH.parent.mkdir(exist_ok=True)

# 임베딩 모델 로드
embedding_model = AutoModel.from_pretrained("BM-K/KoSimCSE-roberta-multitask")
embedding_tokenizer = AutoTokenizer.from_pretrained("BM-K/KoSimCSE-roberta-multitask")
embedding_model.eval()

def extract_text_from_pptx(file_path):
    """PPTX 파일에서 텍스트를 추출합니다."""
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

def extract_text_from_pdf(file_path):
    """PDF 파일에서 텍스트를 추출합니다."""
    pdf_text = ""
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                pdf_text += page_text + "\n"
    return pdf_text

def chunk_text(text, chunk_size=400, chunk_overlap=50):
    """텍스트를 청크로 분할합니다."""
    text = text.strip()
    chunks = []
    start = 0
    text_length = len(text)
    while start < text_length:
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        if end >= text_length:
            break
        start = end - chunk_overlap
    return chunks

@app.get("/manager", response_class=HTMLResponse)
async def data_control_page(request: Request):
    return templates.TemplateResponse("data_manager.html", {"request": request})

@app.post("/upload")
async def data_upload(files: List[UploadFile] = File(...)):
    if not files:
        raise HTTPException(status_code=400, detail="파일이 업로드되지 않았습니다.")
    
    # 기존 데이터 로드
    if DATA_PATH.exists():
        with open(DATA_PATH, "r", encoding="utf-8") as f:
            existing_data = json.load(f)
    else:
        existing_data = []
    
    max_chunk_id = 0
    for entry in existing_data:
        for chunk in entry.get("chunks", []):
            max_chunk_id = max(max_chunk_id, chunk.get("chunk_id", 0))
    
    messages = []
    for file in files:
        try:
            ext = os.path.splitext(file.filename)[1].lower()
            content = ""
            
            if ext == ".txt":
                content = (await file.read()).decode("utf-8")
            elif ext == ".pptx":
                temp_path = DATA_PATH.parent / file.filename
                with open(temp_path, "wb") as f:
                    f.write(await file.read())
                content = extract_text_from_pptx(temp_path)
                temp_path.unlink()
            elif ext == ".pdf":
                temp_path = DATA_PATH.parent / file.filename
                with open(temp_path, "wb") as f:
                    f.write(await file.read())
                content = extract_text_from_pdf(temp_path)
                temp_path.unlink()
            elif ext == ".json":
                content = (await file.read()).decode("utf-8")
                json_data = json.loads(content)
            else:
                messages.append(f"{file.filename}: 지원되지 않는 파일 형식입니다.")
                continue
            
            if ext != ".json":
                if not content.strip():
                    messages.append(f"{file.filename}: 내용이 없습니다.")
                    continue
                
                splitted_contents = chunk_text(content)
                existing_entry = next((entry for entry in existing_data if entry.get("file_name") == file.filename), None)
                if not existing_entry:
                    existing_entry = {
                        "file_name": file.filename,
                        "chunks": []
                    }
                    existing_data.append(existing_entry)
                
                for idx_ch, ch in enumerate(splitted_contents):
                    max_chunk_id += 1
                    vector = vectorize_content(ch)
                    new_chunk = {
                        "chunk_id": max_chunk_id,
                        "title": os.path.splitext(file.filename)[0],
                        "date": datetime.datetime.now().strftime("%Y-%m-%d"),
                        "type": ext.replace(".", ""),
                        "text": ch,
                        "text_short": ch[:200],
                        "vector": vector,
                        "text_vis": normalize_text_vis(ch),
                    }
                    existing_entry["chunks"].append(new_chunk)
                
                messages.append(f"{file.filename}: 업로드 및 벡터화 성공 (총 {len(splitted_contents)}개 청크).")
            else:
                if not isinstance(json_data, list):
                    json_data = [json_data]
                
                for entry_obj in json_data:
                    if "chunks" not in entry_obj:
                        continue
                    for chunk in entry_obj["chunks"]:
                        text = chunk.get("text", "")
                        if not chunk.get("vector"):
                            if text:
                                chunk["vector"] = vectorize_content(text)
                            else:
                                chunk["vector"] = [0.0] * 768
                        if not chunk.get("text_vis"):
                            chunk["text_vis"] = normalize_text_vis(text)
                
                for entry_obj in json_data:
                    for chunk in entry_obj.get("chunks", []):
                        max_chunk_id += 1
                        chunk["chunk_id"] = max_chunk_id
                    existing_data.append(entry_obj)
                
                messages.append(f"{file.filename}: JSON 업로드 및 벡터화/보완 완료.")
        
        except Exception as e:
            messages.append(f"{file.filename}: 업로드 실패: {str(e)}")
    
    try:
        with open(DATA_PATH, "w", encoding="utf-8") as f:
            json.dump(existing_data, f, ensure_ascii=False, indent=2)
    except Exception as save_err:
        raise HTTPException(status_code=500, detail=f"데이터 저장 실패: {str(save_err)}")
    
    return {"message": "\n".join(messages)}

@app.get("/list")
async def data_list(page: int = 1):
    try:
        if DATA_PATH.exists():
            with open(DATA_PATH, "r", encoding="utf-8") as f:
                data_entries = json.load(f)
        else:
            data_entries = []
        
        summary = []
        for idx, entry in enumerate(data_entries):
            if "chunks" in entry and entry["chunks"]:
                chunk = entry["chunks"][0]
                summary.append({
                    "index": idx,
                    "file_name": entry.get("file_name", ""),
                    "title": chunk.get("title", ""),
                    "date": chunk.get("date", "")
                })
            else:
                summary.append({
                    "index": idx,
                    "file_name": entry.get("file_name", ""),
                    "title": "",
                    "date": ""
                })
        
        per_page = 30
        total = len(summary)
        start = (page - 1) * per_page
        end = start + per_page
        paginated = summary[start:end]
        
        return {
            "page": page,
            "per_page": per_page,
            "total": total,
            "data": paginated
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"데이터 목록 불러오기 실패: {str(e)}")

@app.get("/detail/{index}")
async def data_detail(index: int):
    try:
        if not DATA_PATH.exists():
            raise HTTPException(status_code=404, detail="데이터 파일이 존재하지 않습니다.")
        
        with open(DATA_PATH, "r", encoding="utf-8") as f:
            data_entries = json.load(f)
        
        if index < 0 or index >= len(data_entries):
            raise HTTPException(status_code=400, detail="유효하지 않은 인덱스입니다.")
        
        return data_entries[index]
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"데이터 상세 보기 실패: {str(e)}")

@app.post("/delete")
async def data_delete(index: int):
    try:
        if not DATA_PATH.exists():
            raise HTTPException(status_code=404, detail="데이터 파일이 존재하지 않습니다.")
        
        with open(DATA_PATH, "r", encoding="utf-8") as f:
            data_entries = json.load(f)
        
        if not (0 <= index < len(data_entries)):
            raise HTTPException(status_code=400, detail="유효하지 않은 인덱스입니다.")
        
        del data_entries[index]
        
        with open(DATA_PATH, "w", encoding="utf-8") as f:
            json.dump(data_entries, f, ensure_ascii=False, indent=2)
        
        return {"message": "데이터 삭제가 완료되었습니다."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"데이터 삭제 실패: {str(e)}")

@app.get("/search")
async def search_data(q: str):
    if not DATA_PATH.exists():
        return []
    
    with open(DATA_PATH, "r", encoding="utf-8") as f:
        data_entries = json.load(f)
    
    results = []
    for idx, entry in enumerate(data_entries):
        file_name = entry.get("file_name", "").lower()
        title = ""
        if "chunks" in entry and entry["chunks"]:
            title = entry["chunks"][0].get("title", "").lower()
        
        if q.lower() in file_name or q.lower() in title:
            results.append({
                "index": idx,
                "file_name": entry.get("file_name", ""),
                "title": title,
            })
    
    return results 