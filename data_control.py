# data_control.py
import os
import json
import datetime
import torch
import numpy as np
from umap import UMAP
from flask import Blueprint, request, jsonify, render_template
from transformers import AutoModel, AutoTokenizer
from utils import vectorize_content, normalize_text_vis  # Assumes you have defined vectorize_content in utils.py

# For PPTX extraction
from pptx import Presentation

# For PDF extraction
import PyPDF2

# For visualization we use Plotly
import plotly.express as px
import pandas as pd


# Configuration
import yaml
from box import Box
with open("./config.yaml", "r") as f:
    config_yaml = yaml.load(f, Loader=yaml.FullLoader)
    config = Box(config_yaml)

DATA_PATH = config.data_path
print("현재 사용 중인 데이터 : ", DATA_PATH)

# 임베딩 모델 로드 (필요한 경우 캐싱 고려)
embedding_model = AutoModel.from_pretrained("BM-K/KoSimCSE-roberta-multitask")
embedding_tokenizer = AutoTokenizer.from_pretrained("BM-K/KoSimCSE-roberta-multitask")
embedding_model.eval()

data_control_bp = Blueprint("data_manager", __name__, template_folder="templates")

# --- Helper functions for new file types ---

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

def extract_text_from_pdf(file_path):
    pdf_text = ""
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                pdf_text += page_text + "\n"
    return pdf_text

@data_control_bp.route("/manager")
def data_control_page():
    return render_template("data_manager.html")

# --- 다중 파일 업로드 지원 (수정됨) ---
@data_control_bp.route("/upload", methods=["POST"])
def data_upload():
    if "dataFile" not in request.files:
        return jsonify({"message": "파일이 업로드되지 않았습니다."}), 400
    files = request.files.getlist("dataFile")
    if not files:
        return jsonify({"message": "업로드할 파일이 없습니다."}), 400

    # 기존 데이터 로드
    if os.path.exists(DATA_PATH):
        with open(DATA_PATH, "r", encoding="utf-8") as f:
            existing_data = json.load(f)
    else:
        existing_data = []
        
    # 모든 파일에 걸쳐 최대 chunk_id 찾기
    max_chunk_id = 0
    for entry in existing_data:
        for chunk in entry.get("chunks", []):
            max_chunk_id = max(max_chunk_id, chunk.get("chunk_id", 0))
    
    print(f"기존 데이터의 최대 chunk_id: {max_chunk_id}")
    
    messages = []
    for file in files:
        if file.filename == "":
            messages.append("파일 이름이 없습니다.")
            continue
        ext = os.path.splitext(file.filename)[1].lower()
        try:
            new_chunk = None
            # TXT 파일 처리
            if ext == ".txt":
                content = file.read().decode("utf-8")
                vector = vectorize_content(content)
                text_vis = content
                new_chunk = {
                    "chunk_id": None,  # 이후에 할당
                    "title": os.path.splitext(file.filename)[0],
                    "date": datetime.datetime.now().strftime("%Y-%m-%d"),
                    "type": "text",
                    "text": content,
                    "text_short": content[:200],
                    "vector": vector,
                    "text_vis": normalize_text_vis(text_vis)
                }
            # JSON 파일 처리
            elif ext == ".json":
                new_entry = json.load(file)
                if not isinstance(new_entry, list):
                    new_entry = [new_entry]
                # JSON 파일의 각 청크에 대해 벡터 처리
                for entry in new_entry:
                    for chunk in entry.get("chunks", []):
                        if not chunk.get("vector"):
                            text = chunk.get("text", "")
                            chunk["vector"] = vectorize_content(text) if text else [[0.0] * 768]
                        else:
                            vector = chunk["vector"]
                            expected_dim = 768
                            if not isinstance(vector, list):
                                text = chunk.get("text", "")
                                chunk["vector"] = vectorize_content(text) if text else [[0.0]*expected_dim]
                            elif len(vector) != expected_dim and isinstance(vector[0], (int, float)):
                                chunk["vector"] = vector[:expected_dim] if len(vector) >= expected_dim else vector + [0.0]*(expected_dim - len(vector))
                        chunk["text_vis"] = normalize_text_vis(chunk.get("text_vis", ""))
                # JSON 파일은 전체 엔트리 추가
                for entry in new_entry:
                    for chunk in entry.get("chunks", []):
                        max_chunk_id += 1
                        chunk["chunk_id"] = max_chunk_id
                    existing_data.append(entry)
                messages.append(f"{file.filename}: 업로드 및 벡터화 성공.")
                continue  # 다음 파일로 넘어감
            # PPTX 파일 처리
            elif ext == ".pptx":
                temp_path = os.path.join("temp", file.filename)
                os.makedirs("temp", exist_ok=True)
                file.save(temp_path)
                content = extract_text_from_pptx(temp_path)
                os.remove(temp_path)
                vector = vectorize_content(content)
                text_vis = content
                new_chunk = {
                    "chunk_id": None,
                    "title": os.path.splitext(file.filename)[0],
                    "date": datetime.datetime.now().strftime("%Y-%m-%d"),
                    "type": "pptx",
                    "text": content,
                    "text_short": content[:200],
                    "vector": vector,
                    "text_vis": normalize_text_vis(text_vis)
                }
            # PDF 파일 처리
            elif ext == ".pdf":
                temp_path = os.path.join("temp", file.filename)
                os.makedirs("temp", exist_ok=True)
                file.save(temp_path)
                content = extract_text_from_pdf(temp_path)
                os.remove(temp_path)
                vector = vectorize_content(content)
                text_vis = content
                new_chunk = {
                    "chunk_id": None,
                    "title": os.path.splitext(file.filename)[0],
                    "date": datetime.datetime.now().strftime("%Y-%m-%d"),
                    "type": "pdf",
                    "text": content,
                    "text_short": content[:200],
                    "vector": vector,
                    "text_vis": normalize_text_vis(text_vis)
                }
            else:
                messages.append(f"{file.filename}: 지원되지 않는 파일 형식입니다.")
                continue

            # 단일 파일(텍스트, pptx, pdf)의 경우
            if new_chunk is not None:
                existing_entry = next((entry for entry in existing_data if entry.get("file_name") == file.filename), None)
                if existing_entry:
                    max_chunk_id += 1
                    new_chunk["chunk_id"] = max_chunk_id
                    existing_entry.setdefault("chunks", []).append(new_chunk)
                else:
                    max_chunk_id += 1
                    new_chunk["chunk_id"] = max_chunk_id
                    new_entry = {
                        "file_name": file.filename,
                        "chunks": [new_chunk]
                    }
                    existing_data.append(new_entry)
                messages.append(f"{file.filename}: 업로드 및 벡터화 성공.")
        except Exception as e:
            messages.append(f"{file.filename}: 업로드 실패: {str(e)}")
    # 업데이트된 데이터 저장
    with open(DATA_PATH, "w", encoding="utf-8") as f:
        json.dump(existing_data, f, ensure_ascii=False, indent=2)
    return jsonify({"message": "\n".join(messages)})

# --- 페이지네이션 및 인덱스 포함 데이터 목록 (수정됨) ---
@data_control_bp.route("/list", methods=["GET"])
def data_list():
    try:
        if os.path.exists(DATA_PATH):
            with open(DATA_PATH, "r", encoding="utf-8") as f:
                data_entries = json.load(f)
        else:
            data_entries = []
        summary = []
        for idx, entry in enumerate(data_entries):
            if "chunks" in entry and entry["chunks"]:
                chunk = entry["chunks"][0]
                summary.append({
                    "index": idx,
                    "file_name": entry.get("file_name", ""),
                    "title": chunk.get("title", ""),
                    "date": chunk.get("date", "")
                })
            else:
                summary.append({
                    "index": idx,
                    "file_name": entry.get("file_name", ""),
                    "title": "",
                    "date": ""
                })
        # 페이지네이션: 기본 페이지 1, 한 페이지당 30개
        page = request.args.get("page", 1, type=int)
        per_page = 30
        total = len(summary)
        start = (page - 1) * per_page
        end = start + per_page
        paginated = summary[start:end]
        return jsonify({
            "page": page,
            "per_page": per_page,
            "total": total,
            "data": paginated
        })
    except Exception as e:
        return jsonify({"message": f"데이터 목록 불러오기 실패: {str(e)}"}), 500

# --- 상세보기 엔드포인트 추가 ---
@data_control_bp.route("/detail/<int:index>", methods=["GET"])
def data_detail(index):
    try:
        if os.path.exists(DATA_PATH):
            with open(DATA_PATH, "r", encoding="utf-8") as f:
                data_entries = json.load(f)
        else:
            return jsonify({"message": "데이터 파일이 존재하지 않습니다."}), 404
        if index < 0 or index >= len(data_entries):
            return jsonify({"message": "유효하지 않은 인덱스입니다."}), 400
        return jsonify(data_entries[index])
    except Exception as e:
        return jsonify({"message": f"데이터 상세 보기 실패: {str(e)}"}), 500

@data_control_bp.route("/delete", methods=["POST"])
def data_delete():
    try:
        req = request.get_json()
        index = req.get("index")
        if index is None:
            return jsonify({"message": "삭제할 인덱스가 제공되지 않았습니다."}), 400
        if os.path.exists(DATA_PATH):
            with open(DATA_PATH, "r", encoding="utf-8") as f:
                data_entries = json.load(f)
        else:
            return jsonify({"message": "데이터 파일이 존재하지 않습니다."}), 404
        if not (0 <= index < len(data_entries)):
            return jsonify({"message": "유효하지 않은 인덱스입니다."}), 400
        del data_entries[index]
        with open(DATA_PATH, "w", encoding="utf-8") as f:
            json.dump(data_entries, f, ensure_ascii=False, indent=2)
        return jsonify({"message": "데이터 삭제가 완료되었습니다."})
    except Exception as e:
        return jsonify({"message": f"데이터 삭제 실패: {str(e)}"}), 500

# --- 검색 엔드포인트 ---
@data_control_bp.route("/search", methods=["GET"])
def search_data():
    query = request.args.get("q", "").lower()
    if os.path.exists(DATA_PATH):
        with open(DATA_PATH, "r", encoding="utf-8") as f:
            data_entries = json.load(f)
    else:
        data_entries = []
    results = []
    for idx, entry in enumerate(data_entries):
        file_name = entry.get("file_name", "").lower()
        title = ""
        if "chunks" in entry and entry["chunks"]:
            title = entry["chunks"][0].get("title", "").lower()
        if query in file_name or query in title:
            results.append({
                "index": idx,
                "file_name": entry.get("file_name", ""),
                "title": title,
            })
    return jsonify(results)

# --- UMAP 시각화 API (최적화 적용) ---
@data_control_bp.route("/api/umap_data", methods=["GET"])
def get_umap_data():
    try:
        if os.path.exists(DATA_PATH):
            with open(DATA_PATH, "r", encoding="utf-8") as f:
                data_entries = json.load(f)
        else:
            return jsonify({"error": "데이터 파일을 찾을 수 없습니다."}), 404
        
        vectors = []
        metadata = []
        for entry in data_entries:
            for chunk in entry.get("chunks", []):
                if "vector" in chunk and chunk["vector"]:
                    vectors.append(chunk["vector"])
                    metadata.append({
                        "id": f"{entry.get('file_name', 'unknown')}_{chunk.get('chunk_id', '0')}",
                        "file_name": entry.get("file_name", "Unknown"),
                        "title": chunk.get("title", ""),
                        "date": chunk.get("date", ""),
                        "text_short": chunk.get("text_short", "")[:100] + "..."
                    })
        
        if not vectors:
            return jsonify({"error": "벡터 데이터가 없습니다."}), 404

        # Optimization: If the number of vectors exceeds a threshold, downsample them
        MAX_VECTORS = 10000
        if len(vectors) > MAX_VECTORS:
            indices = np.random.choice(len(vectors), MAX_VECTORS, replace=False)
            vectors = [vectors[i] for i in indices]
            metadata = [metadata[i] for i in indices]

        vectors_array = np.array(vectors)
        vectors_array = np.squeeze(vectors_array)
        
        # UMAP embedding
        n_neighbors = min(15, len(vectors) - 1)
        umap_model = UMAP(n_components=2, 
                          n_neighbors=n_neighbors, 
                          min_dist=0.1, 
                          metric='cosine', 
                          random_state=42)
        embedding = umap_model.fit_transform(vectors_array)
        
        nodes = []
        for i, (x, y) in enumerate(embedding):
            nodes.append({
                "id": metadata[i]["id"],
                "x": float(x),
                "y": float(y),
                "file_name": metadata[i]["file_name"],
                "title": metadata[i]["title"],
                "date": metadata[i]["date"],
                "text_short": metadata[i]["text_short"]
            })
        
        # Compute edges (optional, can be heavy for very large datasets)
        edges = []
        if len(vectors) > 1:
            from sklearn.metrics.pairwise import cosine_similarity
            similarity = cosine_similarity(vectors_array)
            threshold = 0.7
            for i in range(len(vectors)):
                for j in range(i+1, len(vectors)):
                    if similarity[i, j] > threshold:
                        edges.append({
                            "source": metadata[i]["id"],
                            "target": metadata[j]["id"],
                            "value": float(similarity[i, j])
                        })
        
        return jsonify({
            "nodes": nodes,
            "edges": edges
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": f"UMAP 시각화 처리 중 오류 발생: {str(e)}"}), 500
