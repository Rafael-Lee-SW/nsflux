# data_control.py

import os
import json
import datetime
import torch
import numpy as np
from umap import UMAP
from flask import Blueprint, request, jsonify, render_template
from transformers import AutoModel, AutoTokenizer
from utils.utils import vectorize_content, normalize_text_vis  # Assumes you have defined vectorize_content in utils.py

# For PPTX extraction
from pptx import Presentation

# For PDF extraction
import PyPDF2

# For visualization we use Plotly
import plotly.express as px
import pandas as pd

# Configuration
import yaml
from box import Box
with open("./config.yaml", "r") as f:
    config_yaml = yaml.load(f, Loader=yaml.FullLoader)
    config = Box(config_yaml)

DATA_PATH = config.data_path
print("현재 사용 중인 데이터 : ", DATA_PATH)

# 임베딩 모델 로드 (필요한 경우 캐싱 고려)
embedding_model = AutoModel.from_pretrained("BM-K/KoSimCSE-roberta-multitask")
embedding_tokenizer = AutoTokenizer.from_pretrained("BM-K/KoSimCSE-roberta-multitask")
embedding_model.eval()

data_control_bp = Blueprint("data_manager", __name__, template_folder="templates")


# --------------------------------------------------------
# HELPER FUNCTIONS
# --------------------------------------------------------

def extract_text_from_pptx(file_path):
    """PPTX 파일에서 텍스트를 추출합니다."""
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

def extract_text_from_pdf(file_path):
    """PDF 파일에서 텍스트를 추출합니다."""
    pdf_text = ""
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                pdf_text += page_text + "\n"
    return pdf_text

def chunk_text(text, chunk_size=400, chunk_overlap=50):
    """
    긴 텍스트를 chunk_size 단위로 분할합니다.
    오버랩(겹침) 길이는 chunk_overlap으로 지정합니다.
    """
    text = text.strip()
    chunks = []
    
    start = 0
    text_length = len(text)
    while start < text_length:
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        
        if end >= text_length:
            break
        start = end - chunk_overlap
    
    return chunks

# --------------------------------------------------------
# ROUTES
# --------------------------------------------------------

@data_control_bp.route("/manager")
def data_control_page():
    return render_template("data_manager.html")


@data_control_bp.route("/upload", methods=["POST"])
def data_upload():
    """
    사용자가 업로드한 (txt/json/pptx/pdf) 파일을 받아서
    데이터를 JSON으로 유지하고, 각 chunk를 벡터화합니다.
    파일이 큰 경우, chunk_text 함수를 이용해 여러 청크로 나눕니다.
    
    print 함수를 통해 업로드/벡터화 프로세스를 상세히 기록합니다.
    """
    if "dataFile" not in request.files:
        print("[Upload] 파일이 업로드되지 않았습니다.")
        return jsonify({"message": "파일이 업로드되지 않았습니다."}), 400
    
    files = request.files.getlist("dataFile")
    if not files:
        print("[Upload] 업로드할 파일이 없습니다.")
        return jsonify({"message": "업로드할 파일이 없습니다."}), 400

    # 기존 데이터 로드
    if os.path.exists(DATA_PATH):
        with open(DATA_PATH, "r", encoding="utf-8") as f:
            existing_data = json.load(f)
    else:
        existing_data = []
        
    # 모든 파일에 걸쳐 현재까지의 최대 chunk_id 찾기
    max_chunk_id = 0
    for entry in existing_data:
        for chunk in entry.get("chunks", []):
            max_chunk_id = max(max_chunk_id, chunk.get("chunk_id", 0))
    
    print(f"[Upload] 기존 데이터의 최대 chunk_id: {max_chunk_id}")
    
    messages = []
    for file in files:
        if file.filename == "":
            messages.append("파일 이름이 없습니다.")
            print("[Upload] 파일 이름이 없습니다.")
            continue
        
        ext = os.path.splitext(file.filename)[1].lower()
        print(f"[Upload] '{file.filename}' 파일 업로드 시작 (확장자: {ext})")
        
        try:
            # 1. 파일 내용을 텍스트로 변환
            content = ""
            if ext == ".txt":
                content = file.read().decode("utf-8")
                print(f"[Upload] TXT 파일 내용 추출 완료 (길이: {len(content)})")
            elif ext == ".pptx":
                temp_path = os.path.join("temp", file.filename)
                os.makedirs("temp", exist_ok=True)
                file.save(temp_path)
                content = extract_text_from_pptx(temp_path)
                os.remove(temp_path)
                print(f"[Upload] PPTX 텍스트 추출 완료 (길이: {len(content)})")
            elif ext == ".pdf":
                temp_path = os.path.join("temp", file.filename)
                os.makedirs("temp", exist_ok=True)
                file.save(temp_path)
                content = extract_text_from_pdf(temp_path)
                os.remove(temp_path)
                print(f"[Upload] PDF 텍스트 추출 완료 (길이: {len(content)})")
            elif ext == ".json":
                print(f"[Upload] JSON 파일, 별도 처리 진행 예정.")
                # JSON 파일은 이후 별도 처리
            else:
                if ext != ".json":
                    msg_unsupported = f"{file.filename}: 지원되지 않는 파일 형식입니다."
                    messages.append(msg_unsupported)
                    print("[Upload]", msg_unsupported)
                continue
            
            # 2. JSON이 아닌 파일(txt/pptx/pdf)은 chunk_text -> 벡터화
            if ext != ".json":
                if not content.strip():
                    msg_no_content = f"{file.filename}: 내용이 없습니다."
                    messages.append(msg_no_content)
                    print("[Upload]", msg_no_content)
                    continue
                
                # chunk_text
                splitted_contents = chunk_text(content, chunk_size=400, chunk_overlap=50)
                print(f"[Upload] 총 {len(splitted_contents)}개 청크로 분할 완료. (파일: {file.filename})")
                
                # 기존 entry 확인
                existing_entry = next((entry for entry in existing_data if entry.get("file_name") == file.filename), None)
                if not existing_entry:
                    existing_entry = {
                        "file_name": file.filename,
                        "chunks": []
                    }
                    existing_data.append(existing_entry)
                
                # 각 청크를 벡터화
                for idx_ch, ch in enumerate(splitted_contents):
                    preview_text = ch[:100].replace("\n", " ")
                    print(f"\n[Upload] 청크 {idx_ch+1}/{len(splitted_contents)}")
                    print(f" - chunk_id 예정: {max_chunk_id + 1}")
                    print(f" - 청크 내용(앞 100자): {preview_text}...")
                    print(f" - 청크 전체 길이: {len(ch)}")
                    
                    max_chunk_id += 1
                    vector = vectorize_content(ch)
                    print(f" - 벡터화 완료 (768차원 여부는 utils 참고)")
                    
                    new_chunk = {
                        "chunk_id": max_chunk_id,
                        "title": os.path.splitext(file.filename)[0],
                        "date": datetime.datetime.now().strftime("%Y-%m-%d"),
                        "type": ext.replace(".", ""),  # "txt", "pptx", "pdf"
                        "text": ch,
                        "text_short": ch[:200],
                        "vector": vector,
                        "text_vis": normalize_text_vis(ch),
                    }
                    existing_entry["chunks"].append(new_chunk)
                
                success_msg = f"{file.filename}: 업로드 및 벡터화 성공 (총 {len(splitted_contents)}개 청크)."
                messages.append(success_msg)
                print("[Upload]", success_msg)
            
            else:
                # 3. JSON 파일 처리
                json_data = json.load(file)
                print(f"[Upload] JSON 로드 완료.")
                
                if not isinstance(json_data, list):
                    json_data = [json_data]
                
                for entry_obj in json_data:
                    if "chunks" not in entry_obj:
                        continue
                    for idx_ck, chunk in enumerate(entry_obj["chunks"]):
                        text = chunk.get("text", "")
                        print(f"\n[Upload][JSON] 청크 {idx_ck+1}/{len(entry_obj['chunks'])} 처리 중.")
                        print(f" - 기존 chunk_id: {chunk.get('chunk_id', 'None')}")
                        print(f" - text 길이: {len(text)}")
                        
                        if not chunk.get("vector"):
                            # vector가 없을 경우 새로 벡터화
                            if text:
                                chunk["vector"] = vectorize_content(text)
                                print(" - 벡터화 완료(신규).")
                            else:
                                # 텍스트가 없다면 0벡터 할당
                                chunk["vector"] = [0.0] * 768
                                print(" - 텍스트 없음 -> 0벡터 할당.")
                        else:
                            # 벡터가 이미 존재한다면 차원 보정
                            vector = chunk["vector"]
                            expected_dim = 768
                            if not isinstance(vector, list):
                                # 잘못된 형식 -> 재벡터화
                                chunk["vector"] = vectorize_content(text) if text else [0.0]*expected_dim
                                print(" - 벡터 형식 이상 -> 재벡터화 처리.")
                            elif len(vector) != expected_dim and isinstance(vector[0], (int, float)):
                                # 길이가 맞지 않는 경우 -> 보정
                                if len(vector) > expected_dim:
                                    chunk["vector"] = vector[:expected_dim]
                                    print(" - 벡터가 너무 길어 잘라냄.")
                                else:
                                    chunk["vector"] = vector + [0.0]*(expected_dim - len(vector))
                                    print(" - 벡터가 짧아 0으로 패딩.")
                            else:
                                print(" - 기존 벡터 유지.")
                        
                        # text_vis 없는 경우 생성
                        if not chunk.get("text_vis"):
                            chunk["text_vis"] = normalize_text_vis(text)
                            print(" - text_vis 필드 생성.")
                
                # JSON 파일에 담긴 entry를 전체 데이터에 추가
                for entry_obj in json_data:
                    # chunk_id 업데이트
                    for chunk in entry_obj.get("chunks", []):
                        max_chunk_id += 1
                        chunk["chunk_id"] = max_chunk_id
                        print(f" - JSON 파일 chunk_id 할당: {max_chunk_id}")
                    existing_data.append(entry_obj)
                
                success_json_msg = f"{file.filename}: JSON 업로드 및 벡터화/보완 완료."
                messages.append(success_json_msg)
                print("[Upload]", success_json_msg)
        
        except Exception as e:
            error_msg = f"{file.filename}: 업로드 실패: {str(e)}"
            messages.append(error_msg)
            print("[Upload]", error_msg)

    # 모든 처리가 끝난 뒤, 업데이트된 데이터 저장
    try:
        with open(DATA_PATH, "w", encoding="utf-8") as f:
            json.dump(existing_data, f, ensure_ascii=False, indent=2)
        print(f"[Upload] 모든 데이터 저장 완료 -> {DATA_PATH}")
    except Exception as save_err:
        print("[Upload] 데이터 저장 중 오류:", str(save_err))
        return jsonify({"message": f"데이터 저장 실패: {str(save_err)}"}), 500
    
    return jsonify({"message": "\n".join(messages)})


@data_control_bp.route("/list", methods=["GET"])
def data_list():
    """
    데이터 목록을 페이지네이션해서 보여줍니다.
    한 페이지에 30개씩.
    """
    try:
        if os.path.exists(DATA_PATH):
            with open(DATA_PATH, "r", encoding="utf-8") as f:
                data_entries = json.load(f)
        else:
            data_entries = []
        
        summary = []
        for idx, entry in enumerate(data_entries):
            if "chunks" in entry and entry["chunks"]:
                chunk = entry["chunks"][0]
                summary.append({
                    "index": idx,
                    "file_name": entry.get("file_name", ""),
                    "title": chunk.get("title", ""),
                    "date": chunk.get("date", "")
                })
            else:
                summary.append({
                    "index": idx,
                    "file_name": entry.get("file_name", ""),
                    "title": "",
                    "date": ""
                })
        
        page = request.args.get("page", 1, type=int)
        per_page = 30
        total = len(summary)
        start = (page - 1) * per_page
        end = start + per_page
        paginated = summary[start:end]
        
        return jsonify({
            "page": page,
            "per_page": per_page,
            "total": total,
            "data": paginated
        })
    except Exception as e:
        return jsonify({"message": f"데이터 목록 불러오기 실패: {str(e)}"}), 500


@data_control_bp.route("/detail/<int:index>", methods=["GET"])
def data_detail(index):
    """
    주어진 index에 해당하는 데이터의 상세 정보를 조회합니다.
    """
    try:
        if os.path.exists(DATA_PATH):
            with open(DATA_PATH, "r", encoding="utf-8") as f:
                data_entries = json.load(f)
        else:
            return jsonify({"message": "데이터 파일이 존재하지 않습니다."}), 404
        
        if index < 0 or index >= len(data_entries):
            return jsonify({"message": "유효하지 않은 인덱스입니다."}), 400
        
        print("[SOOWAN]데이터 상세보기 : ", data_entries[index])
        return jsonify(data_entries[index])
    except Exception as e:
        return jsonify({"message": f"데이터 상세 보기 실패: {str(e)}"}), 500


@data_control_bp.route("/delete", methods=["POST"])
def data_delete():
    """
    주어진 index에 해당하는 엔트리를 삭제합니다.
    """
    try:
        req = request.get_json()
        index = req.get("index")
        if index is None:
            return jsonify({"message": "삭제할 인덱스가 제공되지 않았습니다."}), 400
        
        if os.path.exists(DATA_PATH):
            with open(DATA_PATH, "r", encoding="utf-8") as f:
                data_entries = json.load(f)
        else:
            return jsonify({"message": "데이터 파일이 존재하지 않습니다."}), 404
        
        if not (0 <= index < len(data_entries)):
            return jsonify({"message": "유효하지 않은 인덱스입니다."}), 400
        
        del data_entries[index]
        
        with open(DATA_PATH, "w", encoding="utf-8") as f:
            json.dump(data_entries, f, ensure_ascii=False, indent=2)
        
        return jsonify({"message": "데이터 삭제가 완료되었습니다."})
    except Exception as e:
        return jsonify({"message": f"데이터 삭제 실패: {str(e)}"}), 500


@data_control_bp.route("/search", methods=["GET"])
def search_data():
    """
    파일명 또는 chunks[0].title 기준 간단 검색.
    """
    query = request.args.get("q", "").lower()
    if os.path.exists(DATA_PATH):
        with open(DATA_PATH, "r", encoding="utf-8") as f:
            data_entries = json.load(f)
    else:
        data_entries = []
    
    results = []
    for idx, entry in enumerate(data_entries):
        file_name = entry.get("file_name", "").lower()
        title = ""
        if "chunks" in entry and entry["chunks"]:
            title = entry["chunks"][0].get("title", "").lower()
        
        if query in file_name or query in title:
            results.append({
                "index": idx,
                "file_name": entry.get("file_name", ""),
                "title": title,
            })
    
    return jsonify(results)


@data_control_bp.route("/api/umap_data", methods=["GET"])
def get_umap_data():
    """
    UMAP으로 2차원 임베딩한 뒤,
    nodes, edges 형태로 반환하여 시각화를 지원합니다.
    """
    try:
        if os.path.exists(DATA_PATH):
            with open(DATA_PATH, "r", encoding="utf-8") as f:
                data_entries = json.load(f)
        else:
            return jsonify({"error": "데이터 파일을 찾을 수 없습니다."}), 404
        
        vectors = []
        metadata = []
        for entry in data_entries:
            for chunk in entry.get("chunks", []):
                if "vector" in chunk and chunk["vector"]:
                    vectors.append(chunk["vector"])
                    metadata.append({
                        "id": f"{entry.get('file_name', 'unknown')}_{chunk.get('chunk_id', '0')}",
                        "file_name": entry.get("file_name", "Unknown"),
                        "title": chunk.get("title", ""),
                        "date": chunk.get("date", ""),
                        "text_short": chunk.get("text_short", "")[:100] + "..."
                    })
        
        if not vectors:
            return jsonify({"error": "벡터 데이터가 없습니다."}), 404

        # 너무 많은 벡터를 처리하면 성능 문제가 있을 수 있음 -> downsample
        MAX_VECTORS = 10000
        if len(vectors) > MAX_VECTORS:
            indices = np.random.choice(len(vectors), MAX_VECTORS, replace=False)
            vectors = [vectors[i] for i in indices]
            metadata = [metadata[i] for i in indices]

        vectors_array = np.array(vectors)
        vectors_array = np.squeeze(vectors_array)
        
        # UMAP
        n_neighbors = min(15, len(vectors) - 1)
        umap_model = UMAP(
            n_components=2,
            n_neighbors=n_neighbors,
            min_dist=0.1,
            metric='cosine',
            random_state=42
        )
        embedding = umap_model.fit_transform(vectors_array)
        
        nodes = []
        for i, (x, y) in enumerate(embedding):
            nodes.append({
                "id": metadata[i]["id"],
                "x": float(x),
                "y": float(y),
                "file_name": metadata[i]["file_name"],
                "title": metadata[i]["title"],
                "date": metadata[i]["date"],
                "text_short": metadata[i]["text_short"]
            })
        
        # (Optional) edges 계산 -> 큰 데이터셋에서는 비추천(매우 느림)
        edges = []
        if len(vectors) > 1:
            from sklearn.metrics.pairwise import cosine_similarity
            similarity = cosine_similarity(vectors_array)
            threshold = 0.7
            for i in range(len(vectors)):
                for j in range(i+1, len(vectors)):
                    if similarity[i, j] > threshold:
                        edges.append({
                            "source": metadata[i]["id"],
                            "target": metadata[j]["id"],
                            "value": float(similarity[i, j])
                        })
        
        return jsonify({
            "nodes": nodes,
            "edges": edges
        })
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": f"UMAP 시각화 처리 중 오류 발생: {str(e)}"}), 500
